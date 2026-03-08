import json
from datetime import datetime
from pathlib import Path

from docx import Document
from docx.shared import Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt as PptPt


BASE_DIR = Path(__file__).resolve().parents[1]
OUT_DIR = BASE_DIR / "outputs" / "documents"
OUT_DIR.mkdir(parents=True, exist_ok=True)
PROFILES_PATH = BASE_DIR / "scripts" / "brand_profiles.json"


DEFAULT_PROFILES = [
    {
        "code": "ezworks",
        "company_name": "EZWorks",
        "host": "root.local",
        "primary_color": "1E3A5F",
        "secondary_color": "2563EB",
        "tagline": "Operator-first multi-tenant expense automation",
        "logo_path": "static/logos/ezworks_mark.svg",
    },
    {
        "code": "kioti",
        "company_name": "KIOTI",
        "host": "kioti.local",
        "primary_color": "B91C1C",
        "secondary_color": "F97316",
        "tagline": "Branded expense workflow for employee adoption",
        "logo_path": "static/logos/company_logo.png",
    },
    {
        "code": "lekpartners",
        "company_name": "LEK Partners",
        "host": "lekpartners.local",
        "primary_color": "0F2F57",
        "secondary_color": "2E6DA4",
        "tagline": "Accounting-firm ready tenant operations",
        "logo_path": "",
    },
]


def _hex_to_rgb(value: str) -> tuple[int, int, int]:
    v = (value or "").strip().lstrip("#")
    if len(v) != 6:
        v = "1E3A5F"
    return int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16)


def _load_profiles() -> list[dict]:
    if PROFILES_PATH.exists():
        try:
            raw = json.loads(PROFILES_PATH.read_text(encoding="utf-8"))
            if isinstance(raw, list) and raw:
                return raw
        except Exception:
            pass
    PROFILES_PATH.write_text(json.dumps(DEFAULT_PROFILES, ensure_ascii=False, indent=2), encoding="utf-8")
    return DEFAULT_PROFILES


def _maybe_add_logo(slide, logo_path: Path):
    if not logo_path.exists():
        return
    if logo_path.suffix.lower() not in {".png", ".jpg", ".jpeg"}:
        return
    slide.shapes.add_picture(str(logo_path), Inches(11.0), Inches(0.2), height=Inches(0.45))


def make_pptx(path: Path, profile: dict):
    prs = Presentation()
    primary = _hex_to_rgb(profile.get("primary_color", "1E3A5F"))
    secondary = _hex_to_rgb(profile.get("secondary_color", "2563EB"))
    company = profile.get("company_name", "Company")
    host = profile.get("host", "localhost")
    tagline = profile.get("tagline", "Expense automation")
    logo_path = BASE_DIR / profile.get("logo_path", "")

    def add_brand_bar(slide):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.55))
        bar.fill.solid()
        bar.fill.fore_color.rgb = PptRGB(*primary)
        bar.line.fill.background()
        label = slide.shapes.add_textbox(Inches(0.35), Inches(0.1), Inches(9.5), Inches(0.35))
        tf = label.text_frame
        tf.text = f"{company} | {host}"
        run = tf.paragraphs[0].runs[0]
        run.font.color.rgb = PptRGB(255, 255, 255)
        run.font.bold = True
        run.font.size = PptPt(14)
        _maybe_add_logo(slide, logo_path)

    def title_slide():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_brand_bar(slide)
        title = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.4), Inches(1.2))
        tf = title.text_frame
        tf.text = f"{company} Expense Platform"
        r = tf.paragraphs[0].runs[0]
        r.font.size = PptPt(44)
        r.font.bold = True
        r.font.color.rgb = PptRGB(*primary)
        sub = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(10.8), Inches(1.0))
        sub.text_frame.text = tagline
        sub.text_frame.paragraphs[0].runs[0].font.size = PptPt(24)
        sub.text_frame.paragraphs[0].runs[0].font.color.rgb = PptRGB(*secondary)
        dt = slide.shapes.add_textbox(Inches(0.8), Inches(6.7), Inches(6.0), Inches(0.4))
        dt.text_frame.text = f"Generated: {datetime.now():%Y-%m-%d %H:%M}"

    def bullet_slide(title: str, bullets: list[str]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_brand_bar(slide)
        t = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11), Inches(0.7))
        t.text_frame.text = title
        tr = t.text_frame.paragraphs[0].runs[0]
        tr.font.size = PptPt(30)
        tr.font.bold = True
        tr.font.color.rgb = PptRGB(*primary)
        box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.0), Inches(4.5))
        tf = box.text_frame
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = b
            p.level = 0
            p.font.size = PptPt(20)
            p.font.color.rgb = PptRGB(40, 40, 40)

    title_slide()
    bullet_slide(
        "Operating Model",
        [
            "Root host supports operator workflows and platform governance.",
            "Each customer host uses isolated config and tenant-specific policies.",
            "Tenant admins handle local operations without cross-tenant access.",
        ],
    )
    bullet_slide(
        "Security and Controls",
        [
            "Host-scoped login domain allow-list policy.",
            "Operator-only access for tenant registry and cross-tenant admin controls.",
            "Admin and Platform console split to prevent accidental misuse.",
        ],
    )
    bullet_slide(
        "Immediate Next Steps",
        [
            "Finalize customer onboarding checklist and handoff SOP.",
            "Add audit logging for policy/admin changes.",
            "Run UAT on root + tenant host scenarios with SSO.",
        ],
    )

    prs.save(str(path))


def make_docx(path: Path, profile: dict):
    primary = _hex_to_rgb(profile.get("primary_color", "1E3A5F"))
    company = profile.get("company_name", "Company")
    host = profile.get("host", "localhost")
    tagline = profile.get("tagline", "Expense automation")

    doc = Document()
    heading = doc.add_heading(f"{company} Expense Platform Brief", level=1)
    if heading.runs:
        heading.runs[0].font.color.rgb = RGBColor(*primary)

    intro = doc.add_paragraph(
        f"Host: {host}\n"
        f"Tagline: {tagline}\n\n"
        "This document is an auto-generated branded brief for stakeholder communication."
    )
    intro.runs[0].font.size = Pt(11)

    doc.add_heading("1. Service Structure", level=2)
    for line in [
        "Root host for operator-level platform management.",
        "Tenant host for customer-specific branding and policy.",
        "Role separation: operator admin vs tenant admin.",
    ]:
        doc.add_paragraph(line, style="List Bullet")

    doc.add_heading("2. Implemented Controls", level=2)
    for line in [
        "Per-host config and branding isolation.",
        "Login domain allow-list per host.",
        "Platform-only tenant registry management.",
        "Cross-tenant actions restricted to operator accounts.",
    ]:
        doc.add_paragraph(line, style="List Bullet")

    doc.add_heading("3. UAT Checklist", level=2)
    for line in [
        "Verify login and logout flow on root host.",
        "Verify login policy enforcement on tenant host.",
        "Verify tenant admin can access only tenant admin features.",
        "Verify operator can manage tenant host from platform console.",
    ]:
        doc.add_paragraph(line, style="List Number")

    doc.add_paragraph(f"Generated at: {datetime.now():%Y-%m-%d %H:%M:%S}")
    doc.save(str(path))


def main():
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    profiles = _load_profiles()
    created = []
    for p in profiles:
        code = (p.get("code") or "brand").strip().lower().replace(" ", "_")
        pptx_path = OUT_DIR / f"{code}_Multitenant_Deck_{ts}.pptx"
        docx_path = OUT_DIR / f"{code}_Project_Brief_{ts}.docx"
        make_pptx(pptx_path, p)
        make_docx(docx_path, p)
        created.extend([pptx_path, docx_path])
    for item in created:
        print(item)


if __name__ == "__main__":
    main()
